"""PPTX → Markdown via python-pptx.

We deliberately avoid pandoc for PPTX because PPTX as an *input* format was only
added in pandoc 3.8.3 (Apr 2026), which most users don't have yet. python-pptx
reads the XML directly and has no native dependencies worth mentioning.
"""

from __future__ import annotations

from pathlib import Path

from ..utils import md_escape_cell
from .base import ConversionError, ConvertResult


def _collect_text(shape) -> list[str]:  # noqa: ANN001 - pptx type
    """Return a list of paragraph-level text strings from a shape tree."""
    out: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                out.append(text)
    if getattr(shape, "shapes", None) is not None:
        for sub in shape.shapes:
            out.extend(_collect_text(sub))
    return out


def _table_to_markdown(table) -> list[str]:  # noqa: ANN001 - pptx type
    rows = []
    for row in table.rows:
        cells = [md_escape_cell(cell.text.strip()) for cell in row.cells]
        rows.append(cells)
    if not rows:
        return []
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "|" + "|".join(["---"] * width) + "|"]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return lines


class PptxConverter:
    name = "pptx"
    extensions = (".pptx",)

    def convert(self, src: Path) -> ConvertResult:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ConversionError(
                src, "python-pptx is not installed — reinstall mdpack"
            ) from e

        try:
            prs = Presentation(str(src))
        except Exception as e:
            raise ConversionError(src, f"failed to open pptx: {e}") from e

        lines: list[str] = [f"# {src.stem}\n"]
        non_empty_slides = 0

        for idx, slide in enumerate(prs.slides, start=1):
            heading = None
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                heading = slide.shapes.title.text_frame.text.strip() or None

            slide_lines: list[str] = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_table:
                    table_lines = _table_to_markdown(shape.table)
                    if table_lines:
                        slide_lines.append("")
                        slide_lines.extend(table_lines)
                    continue
                for text in _collect_text(shape):
                    slide_lines.append(text)

            if not heading and not slide_lines:
                continue
            non_empty_slides += 1
            lines.append(f"\n## Slide {idx}" + (f": {heading}" if heading else ""))
            lines.append("")
            for line in slide_lines:
                lines.append(line)

        if non_empty_slides == 0:
            return ConvertResult(body="", title=src.stem, warnings=["no text found in any slide"])

        return ConvertResult(body="\n".join(lines) + "\n", title=src.stem, warnings=[])
