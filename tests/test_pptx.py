from pathlib import Path

import pytest

pptx_mod = pytest.importorskip("pptx")

from mdpack.converters.base import ConversionError  # noqa: E402
from mdpack.converters.pptx_conv import PptxConverter  # noqa: E402


def _make_pptx(path: Path, slides: list[tuple[str, list[str]]]) -> None:
    """Build a small .pptx with (title, bullets) per slide using python-pptx."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content"
    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0] if bullets else ""
        for b in bullets[1:]:
            body.add_paragraph().text = b
    prs.save(str(path))


def test_pptx_basic(tmp_path: Path) -> None:
    src = tmp_path / "deck.pptx"
    _make_pptx(src, [("Introduction", ["First bullet", "Second bullet"])])

    result = PptxConverter().convert(src)

    assert "# deck" in result.body
    assert "## Slide 1: Introduction" in result.body
    assert "First bullet" in result.body
    assert "Second bullet" in result.body


def test_pptx_multiple_slides(tmp_path: Path) -> None:
    src = tmp_path / "multi.pptx"
    _make_pptx(
        src,
        [
            ("Slide A", ["alpha"]),
            ("Slide B", ["beta"]),
        ],
    )
    result = PptxConverter().convert(src)
    assert "## Slide 1: Slide A" in result.body
    assert "## Slide 2: Slide B" in result.body


def test_pptx_missing_source() -> None:
    with pytest.raises(ConversionError):
        PptxConverter().convert(Path("/nonexistent/deck.pptx"))


def test_pptx_module_loaded() -> None:
    """Regression: catch accidental removal of python-pptx from deps."""
    assert pptx_mod is not None
